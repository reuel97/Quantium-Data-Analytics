#!/usr/bin/env python
# coding: utf-8

# In[1]:


import pandas as pd

# Load the data files
transaction_data_path = 'QVI_transaction_data.xlsx'
purchase_behaviour_path = 'QVI_purchase_behaviour.csv'

# Reading the Excel and CSV files
transaction_data = pd.read_excel(transaction_data_path)
purchase_behaviour = pd.read_csv(purchase_behaviour_path)

# Displaying the first few rows of each dataset
transaction_data_head = transaction_data.head()
purchase_behaviour_head = purchase_behaviour.head()

transaction_data_head, purchase_behaviour_head


# In[2]:


# High-level summary of transaction data
transaction_data_summary = transaction_data.describe(include='all')
transaction_data_info = transaction_data.info()

# High-level summary of purchase behaviour data
purchase_behaviour_summary = purchase_behaviour.describe(include='all')
purchase_behaviour_info = purchase_behaviour.info()

# Check for missing values
transaction_data_missing = transaction_data.isnull().sum()
purchase_behaviour_missing = purchase_behaviour.isnull().sum()

transaction_data_summary, transaction_data_info, transaction_data_missing, purchase_behaviour_summary, purchase_behaviour_info, purchase_behaviour_missing


# In[12]:


# Load and clean the data
transaction_data = pd.read_excel(transaction_data_path)
purchase_behaviour = pd.read_csv(purchase_behaviour_path)

# Convert DATE to datetime format
transaction_data['DATE'] = pd.to_datetime(transaction_data['DATE'], origin='1899-12-30', unit='D')

# Detect and handle outliers in PROD_QTY and TOT_SALES
transaction_data_cleaned = transaction_data[(transaction_data['PROD_QTY'] <= 10) & (transaction_data['TOT_SALES'] <= 50)]

# Correct the extraction and conversion of PACK_SIZE
transaction_data_cleaned['PACK_SIZE'] = transaction_data_cleaned['PROD_NAME'].str.extract(r'(\d+)').astype(float)
transaction_data_cleaned['BRAND_NAME'] = transaction_data_cleaned['PROD_NAME'].str.split().str[0]

# Summarize the cleaned data again
cleaned_data_summary_corrected = transaction_data_cleaned.describe(include='all', datetime_is_numeric=True)

# Calculate Average Spending per Transaction
avg_spending_per_transaction = transaction_data_cleaned.groupby('LYLTY_CARD_NBR')['TOT_SALES'].mean().reset_index()
avg_spending_per_transaction.rename(columns={'TOT_SALES': 'AVG_SPENDING'}, inplace=True)

# Calculate Frequency of Purchases
purchase_frequency = transaction_data_cleaned.groupby('LYLTY_CARD_NBR').size().reset_index(name='PURCHASE_FREQUENCY')

# Calculate Popularity of Brands
brand_popularity = transaction_data_cleaned['BRAND_NAME'].value_counts().reset_index()
brand_popularity.columns = ['BRAND_NAME', 'PURCHASE_COUNT']

# Calculate Pack Size Preferences
pack_size_preferences = transaction_data_cleaned['PACK_SIZE'].value_counts().reset_index()
pack_size_preferences.columns = ['PACK_SIZE', 'PURCHASE_COUNT']

# Merging the purchase behaviour data with calculated metrics
customer_metrics = avg_spending_per_transaction.merge(purchase_frequency, on='LYLTY_CARD_NBR')
customer_metrics = customer_metrics.merge(purchase_behaviour, on='LYLTY_CARD_NBR')

# Display summaries and results
(cleaned_data_summary_corrected.head(), 
 avg_spending_per_transaction.head(), 
 purchase_frequency.head(), 
 brand_popularity.head(), 
 pack_size_preferences.head(), 
 customer_metrics.head())


# In[14]:


import pandas as pd

# Load the new data file
data_path = 'QVI_data.csv'
data = pd.read_csv(data_path)

# Display the first few rows of the data
data.head()


# In[15]:


# Convert DATE to datetime format
data['DATE'] = pd.to_datetime(data['DATE'])

# Extract month and year from DATE
data['MONTH'] = data['DATE'].dt.to_period('M')

# Calculate monthly metrics: total sales, number of customers, average transactions per customer
monthly_metrics = data.groupby(['STORE_NBR', 'MONTH']).agg(
    total_sales=('TOT_SALES', 'sum'),
    total_customers=('LYLTY_CARD_NBR', 'nunique'),
    total_transactions=('TXN_ID', 'nunique')
).reset_index()

# Calculate average transactions per customer
monthly_metrics['avg_transactions_per_customer'] = monthly_metrics['total_transactions'] / monthly_metrics['total_customers']

# Display the first few rows of the monthly metrics
monthly_metrics.head()


# In[17]:


from scipy.stats import pearsonr

def calculate_similarity(trial_store, potential_control_store):
    # Merge the metrics for the trial store and the potential control store on month
    merged_metrics = pd.merge(
        monthly_metrics[monthly_metrics['STORE_NBR'] == trial_store],
        monthly_metrics[monthly_metrics['STORE_NBR'] == potential_control_store],
        on='MONTH',
        suffixes=('_trial', '_control')
    )
    
    # Check if there are enough data points for Pearson correlation
    if len(merged_metrics) < 2:
        return float('-inf')  # Return a very low similarity score if there are not enough data points
    
    # Calculate Pearson correlation for total_sales
    sales_corr, _ = pearsonr(merged_metrics['total_sales_trial'], merged_metrics['total_sales_control'])
    
    # Calculate magnitude distance for total_sales
    sales_min = merged_metrics[['total_sales_trial', 'total_sales_control']].min().min()
    sales_max = merged_metrics[['total_sales_trial', 'total_sales_control']].max().max()
    sales_distance = 1 - ((merged_metrics['total_sales_trial'] - merged_metrics['total_sales_control']).abs() - sales_min) / (sales_max - sales_min)
    
    # Calculate the combined similarity score
    similarity_score = sales_corr + sales_distance.mean()
    
    return similarity_score

def find_control_store(trial_store):
    # Get the list of potential control stores (excluding the trial store)
    potential_control_stores = monthly_metrics['STORE_NBR'].unique()
    potential_control_stores = potential_control_stores[potential_control_stores != trial_store]
    
    # Calculate the similarity score for each potential control store
    similarity_scores = {store: calculate_similarity(trial_store, store) for store in potential_control_stores}
    
    # Select the control store with the highest similarity score
    control_store = max(similarity_scores, key=similarity_scores.get)
    
    return control_store

# Find control stores for the trial stores
trial_stores = [77, 86, 88]
control_stores = {trial_store: find_control_store(trial_store) for trial_store in trial_stores}

control_stores


# In[18]:


# Check the available periods for the trial and control stores
available_periods = monthly_metrics.groupby('STORE_NBR')['MONTH'].agg(['min', 'max']).reset_index()

# Display the available periods for the trial and control stores
trial_and_control_stores = trial_stores + list(control_stores.values())
available_periods = available_periods[available_periods['STORE_NBR'].isin(trial_and_control_stores)]

available_periods


# In[20]:


def compare_stores(trial_store, control_store):
    trial_data = trial_period_data[trial_period_data['STORE_NBR'] == trial_store]
    control_data = trial_period_data[trial_period_data['STORE_NBR'] == control_store]
    
    # Merge trial and control data on month
    merged_data = pd.merge(trial_data, control_data, on='MONTH', suffixes=('_trial', '_control'))
    
    # Calculate the differences in metrics
    merged_data['sales_diff'] = merged_data['total_sales_trial'] - merged_data['total_sales_control']
    merged_data['customers_diff'] = merged_data['total_customers_trial'] - merged_data['total_customers_control']
    merged_data['transactions_diff'] = merged_data['avg_transactions_per_customer_trial'] - merged_data['avg_transactions_per_customer_control']
    
    return merged_data

# Redefine the trial period
trial_period_start = pd.Period('2019-01')
trial_period_end = pd.Period('2019-06')

# Filter data for the updated trial period
trial_period_data = monthly_metrics[(monthly_metrics['MONTH'] >= trial_period_start) & (monthly_metrics['MONTH'] <= trial_period_end)]

# Compare each trial store with its corresponding control store
comparisons = {trial_store: compare_stores(trial_store, control_stores[trial_store]) for trial_store in trial_stores}

# Display the comparison results for one of the trial stores
comparisons[77]


# In[21]:


from scipy.stats import ttest_rel

# Perform paired t-test for total sales between trial and control stores
def test_significance(trial_store, control_store):
    trial_data = trial_period_data[trial_period_data['STORE_NBR'] == trial_store]['total_sales']
    control_data = trial_period_data[trial_period_data['STORE_NBR'] == control_store]['total_sales']
    
    t_stat, p_value = ttest_rel(trial_data, control_data)
    return t_stat, p_value

# Test significance for each trial store
significance_tests = {trial_store: test_significance(trial_store, control_stores[trial_store]) for trial_store in trial_stores}

# Display the results of the significance tests
significance_tests


# In[23]:


import matplotlib.pyplot as plt

# Create a function to plot the comparisons
def plot_comparison(data, metric, title, ylabel):
    plt.figure(figsize=(10, 6))
    
    for store_type in ['trial', 'control']:
        plt.plot(data['MONTH'], data[f'{metric}_{store_type}'], label=f'{store_type.capitalize()} Store', marker='o')
    
    plt.title(title)
    plt.xlabel('Month')
    plt.ylabel(ylabel)
    plt.legend()
    plt.grid(True)
    plt.tight_layout()

    # Save the figure
    plt.savefig(f'{title}.png')
    plt.show()

# Convert the MONTH column to string format suitable for plotting
comparisons[77]['MONTH'] = comparisons[77]['MONTH'].astype(str)

# Plot total sales comparison for Trial Store 77
plot_comparison(comparisons[77], 'total_sales', 'Total Sales Comparison - Store 77', 'Total Sales')

# Plot number of customers comparison for Trial Store 77
plot_comparison(comparisons[77], 'total_customers', 'Number of Customers Comparison - Store 77', 'Number of Customers')

# Plot average transactions per customer comparison for Trial Store 77
plot_comparison(comparisons[77], 'avg_transactions_per_customer', 'Average Transactions per Customer Comparison - Store 77', 'Average Transactions per Customer')


# In[24]:


# Convert the MONTH column to string format for plotting
comparisons[86]['MONTH'] = comparisons[86]['MONTH'].astype(str)
comparisons[88]['MONTH'] = comparisons[88]['MONTH'].astype(str)

# Plot total sales comparison for Trial Store 86
plot_comparison(comparisons[86], 'total_sales', 'Total Sales Comparison - Store 86', 'Total Sales')

# Plot number of customers comparison for Trial Store 86
plot_comparison(comparisons[86], 'total_customers', 'Number of Customers Comparison - Store 86', 'Number of Customers')

# Plot average transactions per customer comparison for Trial Store 86
plot_comparison(comparisons[86], 'avg_transactions_per_customer', 'Average Transactions per Customer Comparison - Store 86', 'Average Transactions per Customer')

# Plot total sales comparison for Trial Store 88
plot_comparison(comparisons[88], 'total_sales', 'Total Sales Comparison - Store 88', 'Total Sales')

# Plot number of customers comparison for Trial Store 88
plot_comparison(comparisons[88], 'total_customers', 'Number of Customers Comparison - Store 88', 'Number of Customers')

# Plot average transactions per customer comparison for Trial Store 88
plot_comparison(comparisons[88], 'avg_transactions_per_customer', 'Average Transactions per Customer Comparison - Store 88', 'Average Transactions per Customer')


# In[28]:


from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide Layouts
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Quantium Project Report"
subtitle.text = "Evaluation of Store Trial Performance"

# Add a slide for each trial store
trial_stores_info = {
    77: "Total Sales, Number of Customers, and Average Transactions per Customer Comparison for Store 77",
    86: "Total Sales, Number of Customers, and Average Transactions per Customer Comparison for Store 86",
    88: "Total Sales, Number of Customers, and Average Transactions per Customer Comparison for Store 88"
}

for trial_store, description in trial_stores_info.items():
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = f"Store {trial_store} Analysis"
    content.text = description
    
    # Add images to the slide
    for metric in ['Total Sales', 'Number of Customers', 'Average Transactions per Customer']:
        img_path = f'{metric} Comparison - Store {trial_store}.png'
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(3.5))

# Recommendations Slide
slide = prs.slides.add_slide(content_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Recommendations"
content.text = (
    "1. **Reevaluate Trial Strategy**: Consider revising the trial strategy, such as extending the trial period, modifying promotional tactics, "
    "or selecting different trial stores to achieve a more noticeable impact.\n\n"
    "2. **Further Analysis**: Conduct a more detailed analysis to identify any potential trends or patterns that might not be captured by the overall sales metrics. "
    "This could include customer feedback, product-specific performance, or regional differences.\n\n"
    "3. **Targeted Promotions**: Implement targeted promotions for high-frequency purchasers to increase loyalty and spending.\n\n"
    "4. **Stock Optimization**: Focus on stocking the most popular pack sizes (175g, 150g) to ensure availability and meet customer demand.\n\n"
    "5. **Brand Partnership Opportunities**: Leverage the popularity of top brands to negotiate better deals or co-promotional opportunities.\n\n"
    "6. **Customer Segmentation**: Use the purchase behaviour data to segment customers based on their lifecycle stage and premium status. "
    "Tailor marketing campaigns to each segment.\n\n"
    "7. **Cross-Selling and Upselling**: Introduce complementary products in the same category or related categories (e.g., dips, sodas) to increase basket size. "
    "Highlight these products during checkout or through personalized recommendations."
)

# Save the presentation
prs.save('Quantium_Project_Report.pptx')

